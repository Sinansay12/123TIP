"""
Document Parsing Service.
Extracts text content from PDF and PPTX files page by page.
"""
import os
from typing import List, Tuple
from PyPDF2 import PdfReader
from pptx import Presentation


class DocumentParser:
    """Handles parsing of PDF and PPTX files."""
    
    @staticmethod
    def parse_pdf(file_path: str) -> List[Tuple[int, str]]:
        """
        Parse PDF and extract text per page.
        
        Args:
            file_path: Path to the PDF file.
            
        Returns:
            List of tuples: (page_number, text_content)
        """
        pages = []
        reader = PdfReader(file_path)
        
        for i, page in enumerate(reader.pages, start=1):
            text = page.extract_text()
            if text and text.strip():
                pages.append((i, text.strip()))
        
        return pages
    
    @staticmethod
    def parse_pptx(file_path: str) -> List[Tuple[int, str]]:
        """
        Parse PowerPoint and extract text per slide.
        
        Args:
            file_path: Path to the PPTX file.
            
        Returns:
            List of tuples: (slide_number, text_content)
        """
        slides = []
        prs = Presentation(file_path)
        
        for i, slide in enumerate(prs.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            
            combined_text = "\n".join(slide_text).strip()
            if combined_text:
                slides.append((i, combined_text))
        
        return slides
    
    def parse(self, file_path: str) -> List[Tuple[int, str]]:
        """
        Parse document based on file extension.
        
        Args:
            file_path: Path to the document.
            
        Returns:
            List of tuples: (page/slide_number, text_content)
        """
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext == ".pdf":
            return self.parse_pdf(file_path)
        elif ext in [".pptx", ".ppt"]:
            return self.parse_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file type: {ext}")
