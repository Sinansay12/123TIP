"""
Content Import Script
Extracts slide content from PPTX/PDF files and imports into database
Uses synchronous SQLAlchemy for standalone script execution
"""
import os
import sys
import io
from pathlib import Path
from typing import List, Dict, Optional

# Fix console encoding for Windows
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')

def safe_print(msg):
    """Print with fallback for encoding issues"""
    try:
        print(msg)
    except UnicodeEncodeError:
        print(msg.encode('ascii', 'replace').decode('ascii'))

# Add parent to path
sys.path.insert(0, str(Path(__file__).parent.parent))

from sqlalchemy import create_engine
from sqlalchemy.orm import sessionmaker, declarative_base
from datetime import datetime

# Try to import pptx for PowerPoint parsing
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("Warning: python-pptx not installed. Install with: pip install python-pptx")

# Database setup - using sync engine for script
DATABASE_URL = "sqlite:///./medical_app.db"  # Main application database
engine = create_engine(DATABASE_URL, echo=False)
SessionLocal = sessionmaker(bind=engine)
Base = declarative_base()

# Define Slide model locally for script
from sqlalchemy import Column, Integer, String, Text, DateTime, JSON

class Slide(Base):
    __tablename__ = "slides"
    
    id = Column(Integer, primary_key=True, index=True)
    department = Column(String(255), nullable=False, index=True)
    topic = Column(String(500), nullable=False)
    page_number = Column(Integer, nullable=False)
    title = Column(String(500), nullable=True)
    content = Column(Text, nullable=False)
    bullet_points = Column(JSON, nullable=True)
    image_url = Column(String(1000), nullable=True)
    professor = Column(String(255), nullable=True)
    source_file = Column(String(500), nullable=True)
    created_at = Column(DateTime, default=datetime.utcnow)


# Department mapping from folder names
DEPARTMENT_MAPPING = {
    "Adli Tıp": "Adli Tıp",
    "Anestezi ve Reanimasyon": "Anestezi ve Reanimasyon",
    "Dermatoloji": "Dermatoloji",
    "Fizik Tedavi ve Rehabilitasyon": "Fizik Tedavi ve Rehabilitasyon",
    "Göz": "Göz",
    "Göğüs Cerrahisi": "Göğüs Cerrahisi",
    "Göğüs Hastalıkları": "Göğüs Hastalıkları",
    "Halk Sağlığı": "Halk Sağlığı",
    "Kalp ve Damar Cerrahisi": "Kalp ve Damar Cerrahisi",
    "Nöroşiruji": "Nöroşiruji",
    "Plastik": "Plastik",
    "Psikiyatri": "Psikiyatri",
    "kardiyoloji": "Kardiyoloji",
    "kulak burun boğaz": "Kulak Burun Boğaz",
    "nöroloji": "Nöroloji",
    "ortopedi ve travmatoloji": "Ortopedi ve Travmatoloji",
}


def extract_professor_from_path(path: str) -> Optional[str]:
    """Extract professor name from folder path"""
    path_lower = path.lower()
    if "prof.dr.hakan kar" in path_lower:
        return "Prof. Dr. Hakan Kar"
    elif "prof.dr.halis dokgöz" in path_lower or "halis" in path_lower:
        return "Prof. Dr. Halis Dokgöz"
    elif "prof.dr.nursel" in path_lower:
        return "Prof. Dr. Nursel G. Bilgin"
    elif "prof.dr.özlem" in path_lower:
        return "Prof. Dr. Özlem Bölgen Çimen"
    elif "ali biçer" in path_lower:
        return "Prof. Dr. Ali Biçer"
    elif "orhan güvener" in path_lower:
        return "Orhan Güvener"
    return None


def extract_topic_from_filename(filename: str) -> str:
    """Extract topic name from filename"""
    name = Path(filename).stem
    name = name.replace("_", " ").replace("-", " ")
    for pattern in ["son", "yeni", "(1)", "(2)", "adlı dosyanın kopyası"]:
        name = name.replace(pattern, "")
    return name.strip()


def extract_slides_from_pptx(file_path: str) -> List[Dict]:
    """Extract slide content from a PPTX file"""
    if not HAS_PPTX:
        return []
    
    slides_data = []
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            title = ""
            content_parts = []
            bullet_points = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if not text:
                        continue
                    
                    if not title and shape.has_text_frame:
                        if shape.text_frame.paragraphs:
                            title = shape.text_frame.paragraphs[0].text.strip()
                            for para in shape.text_frame.paragraphs[1:]:
                                if para.text.strip():
                                    bullet_points.append(para.text.strip())
                    else:
                        content_parts.append(text)
            
            content = "\n".join(content_parts) if content_parts else title
            
            slides_data.append({
                "page_number": i + 1,
                "title": title if title else None,
                "content": content if content else "İçerik yükleniyor...",
                "bullet_points": bullet_points if bullet_points else None,
            })
    except Exception as e:
        safe_print(f"  Error processing {file_path}: {e}")
    
    return slides_data


def import_department_slides(base_path: str, department: str, db) -> int:
    """Import all slides for a department"""
    imported = 0
    slide_dirs = ["SLAYT", "SLAYTLAR", "GÜNCELLENEN SLAYTLAR"]
    
    dept_path = Path(base_path) / department
    if not dept_path.exists():
        return 0
    
    for slide_dir in slide_dirs:
        slides_path = dept_path / slide_dir
        if not slides_path.exists():
            continue
        
        for file_path in slides_path.rglob("*.ppt*"):
            topic = extract_topic_from_filename(file_path.name)
            professor = extract_professor_from_path(str(file_path))
            
            safe_print(f"  [FILE] {topic[:50]}...")
            
            slides_data = extract_slides_from_pptx(str(file_path))
            
            if not slides_data:
                slides_data = [{
                    "page_number": 1,
                    "title": topic,
                    "content": f"{topic}",
                    "bullet_points": None,
                }]
            
            for slide_data in slides_data:
                slide = Slide(
                    department=DEPARTMENT_MAPPING.get(department, department),
                    topic=topic,
                    page_number=slide_data["page_number"],
                    title=slide_data.get("title"),
                    content=slide_data["content"],
                    bullet_points=slide_data.get("bullet_points"),
                    professor=professor,
                    source_file=file_path.name,
                )
                db.add(slide)
                imported += 1
    
    return imported


def main():
    """Main import function"""
    source_path = r"C:\Users\Lenovo\OneDrive\Desktop\tıp5. sınıf"
    
    if not Path(source_path).exists():
        safe_print(f"[ERROR] Kaynak klasor bulunamadi: {source_path}")
        return
    
    safe_print("[INFO] Veritabani olusturuluyor...")
    Base.metadata.create_all(bind=engine)
    
    db = SessionLocal()
    
    try:
        total_imported = 0
        
        for folder in sorted(Path(source_path).iterdir()):
            if folder.is_dir():
                dept_name = folder.name
                safe_print(f"\n[FOLDER] {dept_name}")
                count = import_department_slides(source_path, dept_name, db)
                total_imported += count
                if count > 0:
                    safe_print(f"  [OK] {count} slayt eklendi")
        
        db.commit()
        safe_print(f"\n[SUCCESS] Toplam: {total_imported} slayt ice aktarildi!")
        safe_print(f"[DB] Veritabani: {DATABASE_URL}")
        
    except Exception as e:
        safe_print(f"[ERROR] Hata: {e}")
        db.rollback()
    finally:
        db.close()


if __name__ == "__main__":
    main()
